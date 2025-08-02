import aiofiles
import os
from typing import Tuple
from fastapi import UploadFile
import PyPDF2
import docx
import openpyxl
from pptx import Presentation
import requests
from bs4 import BeautifulSoup
from youtube_transcript_api import YouTubeTranscriptApi
import re
from app.core.config import settings

class DocumentProcessor:
    
    async def process_file(self, file: UploadFile) -> Tuple[str, str]:
        """Process uploaded file and return content and file type"""
        
        # Read file content
        content = await file.read()
        filename = file.filename.lower()
        
        if filename.endswith('.pdf'):
            return await self._process_pdf(content), "PDF"
        elif filename.endswith('.docx'):
            return await self._process_docx(content), "DOCX"
        elif filename.endswith('.xlsx') or filename.endswith('.xls'):
            return await self._process_excel(content), "EXCEL"
        elif filename.endswith('.pptx'):
            return await self._process_pptx(content), "PPT"
        elif filename.endswith('.txt'):
            return content.decode('utf-8'), "TXT"
        elif filename.endswith('.md'):
            return content.decode('utf-8'), "MD"
        else:
            # Try to treat as text
            try:
                return content.decode('utf-8'), "TXT"
            except UnicodeDecodeError:
                raise ValueError(f"Unsupported file type: {filename}")
    
    async def process_url(self, url: str) -> Tuple[str, str]:
        """Process URL and return content"""
        if 'youtube.com' in url or 'youtu.be' in url:
            return await self._process_youtube(url), "youtube"
        else:
            return await self._process_web_url(url), "URL"
    
    async def _process_pdf(self, content: bytes) -> str:
        """Extract text from PDF"""
        try:
            # Save temporarily
            temp_file = f"temp_{os.urandom(8).hex()}.pdf"
            async with aiofiles.open(temp_file, 'wb') as f:
                await f.write(content)
            
            # Extract text
            text = ""
            with open(temp_file, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
            
            # Clean up
            os.remove(temp_file)
            return text.strip()
            
        except Exception as e:
            raise ValueError(f"Failed to process PDF: {str(e)}")
    
    async def _process_docx(self, content: bytes) -> str:
        """Extract text from DOCX"""
        try:
            # Save temporarily
            temp_file = f"temp_{os.urandom(8).hex()}.docx"
            async with aiofiles.open(temp_file, 'wb') as f:
                await f.write(content)
            
            # Extract text
            doc = docx.Document(temp_file)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            
            # Clean up
            os.remove(temp_file)
            return text.strip()
            
        except Exception as e:
            raise ValueError(f"Failed to process DOCX: {str(e)}")
    
    async def _process_excel(self, content: bytes) -> str:
        """Extract text from Excel"""
        try:
            # Save temporarily
            temp_file = f"temp_{os.urandom(8).hex()}.xlsx"
            async with aiofiles.open(temp_file, 'wb') as f:
                await f.write(content)
            
            # Extract text
            workbook = openpyxl.load_workbook(temp_file)
            text = ""
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"Sheet: {sheet_name}\n"
                
                for row in sheet.iter_rows():
                    row_text = []
                    for cell in row:
                        if cell.value is not None:
                            row_text.append(str(cell.value))
                    if row_text:
                        text += "\t".join(row_text) + "\n"
                text += "\n"
            
            # Clean up
            os.remove(temp_file)
            return text.strip()
            
        except Exception as e:
            raise ValueError(f"Failed to process Excel: {str(e)}")
    
    async def _process_pptx(self, content: bytes) -> str:
        """Extract text from PowerPoint"""
        try:
            # Save temporarily
            temp_file = f"temp_{os.urandom(8).hex()}.pptx"
            async with aiofiles.open(temp_file, 'wb') as f:
                await f.write(content)
            
            # Extract text
            prs = Presentation(temp_file)
            text = ""
            
            for i, slide in enumerate(prs.slides):
                text += f"Slide {i + 1}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
            
            # Clean up
            os.remove(temp_file)
            return text.strip()
            
        except Exception as e:
            raise ValueError(f"Failed to process PowerPoint: {str(e)}")
    
    async def _process_web_url(self, url: str) -> str:
        """Extract text from web URL"""
        try:
            response = requests.get(url, timeout=30)
            response.raise_for_status()
            
            soup = BeautifulSoup(response.content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text
            text = soup.get_text()
            
            # Clean up text
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)
            
            return text
            
        except Exception as e:
            raise ValueError(f"Failed to process URL: {str(e)}")
    
    async def _process_youtube(self, url: str) -> str:
        """Extract transcript from YouTube video"""
        try:
            # Extract video ID
            video_id = None
            if 'youtube.com' in url:
                video_id = re.search(r'v=([^&]+)', url)
                if video_id:
                    video_id = video_id.group(1)
            elif 'youtu.be' in url:
                video_id = url.split('/')[-1]
            
            if not video_id:
                raise ValueError("Could not extract video ID from URL")
            
            # Get transcript
            transcript = YouTubeTranscriptApi.get_transcript(video_id)
            
            # Combine transcript text
            text = " ".join([entry['text'] for entry in transcript])
            
            return text
            
        except Exception as e:
            raise ValueError(f"Failed to process YouTube video: {str(e)}")